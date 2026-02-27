#!/usr/bin/env python3
"""Generate Contract Management Process Flow PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
BLUE = RGBColor(0x1A, 0x3A, 0x6B)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
GOLD = RGBColor(0xE6, 0xA8, 0x17)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GREY = RGBColor(0x6B, 0x7C, 0x8D)
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF5, 0xFB)
LIGHT_GREEN_BG = RGBColor(0xEA, 0xFA, 0xF1)
LIGHT_YELLOW_BG = RGBColor(0xFE, 0xF9, 0xE7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.width = border_width
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, NAVY)
    add_text(slide, Inches(1), Inches(2), Inches(11.3), Inches(1.5),
             title, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8),
             subtitle, size=18, bold=False, color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
             "February 2026  •  Confidential", size=12, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def add_section_slide(title, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(0.5), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = str(num)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_text(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    return slide


def add_step_box(slide, left, top, width, height, step_num, title, description, color=ACCENT):
    shape = add_shape(slide, left, top, width, height, fill_color=LIGHT_BLUE_BG, border_color=color)
    # Step number
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left + Inches(0.15), top + Inches(0.15),
                                   Inches(0.35), Inches(0.35))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = str(step_num)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_text(slide, left + Inches(0.6), top + Inches(0.1), width - Inches(0.75), Inches(0.35),
             title, size=13, bold=True, color=NAVY)
    # Description
    add_text(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
             description, size=10, color=GREY)


def add_arrow(slide, start_left, start_top, end_left, end_top, color=ACCENT):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
add_title_slide(
    "Contract Management Process Flows",
    "Trammo Trading Desk — 7-Step Contract Wizard & Solver Integration"
)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Overview
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Contract Management Overview", 1)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
         "The Contract Management module is a guided 7-step wizard that takes traders from product group selection\n"
         "through counterparty identification, commercial negotiation, clause selection, LP solver validation,\n"
         "review, and final approval — all within a single real-time LiveView interface.",
         size=14, color=GREY)

# 7 step boxes
step_data = [
    (1, "Product Group", "Select NH3 Domestic,\nSulphur Intl, Petcoke,\nor NH3 International"),
    (2, "Counterparty", "Name, commodity,\ndelivery window"),
    (3, "Commercial", "Quantity, price,\nfreight, payment,\nincoterm"),
    (4, "Clauses", "Select applicable\ncontract clauses"),
    (5, "Optimizer", "HiGHS LP Solver\n+ Claude AI analysis"),
    (6, "Review", "Full summary of\nall entered data"),
    (7, "Approval", "Submit & create\nDB records"),
]
x_start = Inches(0.5)
box_w = Inches(1.65)
box_h = Inches(1.6)
gap = Inches(0.18)
for i, (num, title, desc) in enumerate(step_data):
    left = x_start + i * (box_w + gap)
    color = GREEN if num in (5, 7) else ACCENT
    add_step_box(slide, left, Inches(2.5), box_w, box_h, num, title, desc, color)

# Key differentiator callout
callout = add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.2),
                    fill_color=LIGHT_GREEN_BG, border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(1.0), Inches(4.6), Inches(11.1), Inches(1.0),
         "KEY DIFFERENTIATOR: Unlike traditional contract management systems, Trammo integrates a linear-programming\n"
         "solver directly into Step 5 — giving traders quantitative validation of contract economics\n"
         "BEFORE approval, not after. Claude AI then explains the results in natural language.",
         size=12, bold=False, color=DARK_TEXT)

# Footer
add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: End-to-End Flow
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("End-to-End Contract Flow", 2)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "From the moment a trader starts a new contract to the point it becomes an active constraint in the solver.",
         size=14, color=GREY)

# Flow boxes
flow_items = [
    ("Step 1", "Product Group\nSelection", ACCENT, Inches(0.5), Inches(2.2)),
    ("Step 2", "Counterparty &\nCommodity", ACCENT, Inches(2.2), Inches(2.2)),
    ("Step 3", "Commercial\nTerms", ACCENT, Inches(3.9), Inches(2.2)),
    ("Step 4", "Clause\nSelection", ACCENT, Inches(5.6), Inches(2.2)),
    ("Step 5", "Optimizer\nValidation", GREEN, Inches(7.3), Inches(2.2)),
    ("Step 6", "Review\nSummary", ACCENT, Inches(9.0), Inches(2.2)),
    ("Step 7", "Approval\nSubmission", GREEN, Inches(10.7), Inches(2.2)),
]
for label, text, color, left, top in flow_items:
    box = add_shape(slide, left, top, Inches(1.5), Inches(1.1), fill_color=WHITE, border_color=color, border_width=Pt(2))
    add_text(slide, left + Inches(0.05), top + Inches(0.05), Inches(1.4), Inches(0.25),
             label, size=9, bold=True, color=color)
    add_text(slide, left + Inches(0.05), top + Inches(0.3), Inches(1.4), Inches(0.7),
             text, size=11, bold=False, color=NAVY)

# Lower section: what happens after approval
after_items = [
    ("DB Transaction", "Negotiation +\nContract + Version\n+ Approval records", ACCENT, Inches(1.5), Inches(4.0)),
    ("Pending\nApproval", "Trading Manager\nreviews", ORANGE, Inches(4.0), Inches(4.0)),
    ("Active\nContract", "Approved &\noperational", GREEN, Inches(6.5), Inches(4.0)),
    ("Solver\nConstraint", "Contract becomes\nlive constraint", PURPLE, Inches(9.0), Inches(4.0)),
]
for label, text, color, left, top in after_items:
    box = add_shape(slide, left, top, Inches(2.0), Inches(1.3), fill_color=WHITE, border_color=color, border_width=Pt(2))
    add_text(slide, left + Inches(0.1), top + Inches(0.08), Inches(1.8), Inches(0.4),
             label, size=11, bold=True, color=color)
    add_text(slide, left + Inches(0.1), top + Inches(0.55), Inches(1.8), Inches(0.7),
             text, size=10, color=GREY)

add_text(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.6),
         "Every step transition emits an event to the EventEmitter pipeline → vectorized → stored in pgvector for semantic search.",
         size=12, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Steps 1-3 Detail
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Steps 1-3: Product, Counterparty, Commercial Terms", 3)

# Step 1
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(4.5), fill_color=WHITE, border_color=ACCENT)
add_text(slide, Inches(0.7), Inches(1.6), Inches(3.4), Inches(0.4),
         "Step 1: Product Group Selection", size=14, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(2.1), Inches(3.4), Inches(3.5),
         "Four product groups available:\n\n"
         "• NH3 Domestic Barge (AD)\n"
         "  Anhydrous Ammonia, Aqua Ammonia\n\n"
         "• Sulphur International (SI)\n"
         "  Granular, Liquid, Formed\n\n"
         "• Petcoke (PC)\n"
         "  Fuel-Grade, Anode-Grade, Calcined\n\n"
         "• NH3 International (AI)\n"
         "  Anhydrous Ammonia, Ammonia Solution\n\n"
         "Validation: Must select one to proceed.",
         size=11, color=DARK_TEXT)

# Step 2
add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(4.5), fill_color=WHITE, border_color=ACCENT)
add_text(slide, Inches(4.8), Inches(1.6), Inches(3.4), Inches(0.4),
         "Step 2: Counterparty & Commodity", size=14, bold=True, color=NAVY)
add_text(slide, Inches(4.8), Inches(2.1), Inches(3.4), Inches(3.5),
         "Fields collected:\n\n"
         "• Counterparty Name (required)\n"
         "  e.g. ACME Trading Corp\n\n"
         "• Commodity (required)\n"
         "  Filtered by product group\n\n"
         "• Delivery Window Start (required)\n"
         "  Date picker\n\n"
         "• Delivery Window End (required)\n"
         "  Date picker\n\n"
         "Validation: All four fields required.",
         size=11, color=DARK_TEXT)

# Step 3
add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(4.5), fill_color=WHITE, border_color=ACCENT)
add_text(slide, Inches(8.9), Inches(1.6), Inches(3.7), Inches(0.4),
         "Step 3: Commercial Terms", size=14, bold=True, color=NAVY)
add_text(slide, Inches(8.9), Inches(2.1), Inches(3.7), Inches(3.5),
         "Fields collected:\n\n"
         "• Quantity + Unit (MT/KT/BBL) — required\n\n"
         "• Proposed Price (USD) — required\n"
         "  → Maps to solver var: nh3_price\n\n"
         "• Proposed Freight (USD/MT) — optional\n"
         "  → Maps to solver var: barge_freight\n\n"
         "• Payment Terms — default: Net 30\n"
         "  Options: Net 30/45/60, LC, Prepay\n\n"
         "• Incoterm — default: FOB\n"
         "  Options: FOB, CFR, CIF, DAP, etc.\n\n"
         "Solver integration: Price & freight override\n"
         "product group defaults in Step 5.",
         size=11, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: Steps 4-5 Detail
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Steps 4-5: Clause Selection & Optimizer Validation", 4)

# Step 4
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5), fill_color=WHITE, border_color=ACCENT)
add_text(slide, Inches(0.7), Inches(1.6), Inches(5.1), Inches(0.4),
         "Step 4: Clause Selection", size=14, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(2.1), Inches(5.1), Inches(3.5),
         "10 standard clauses — toggle on/off:\n\n"
         "  ✓ Force Majeure (pre-selected)\n"
         "  • Demurrage & Dispatch\n"
         "  • Quality Specification\n"
         "  • Quantity Tolerance (±5%)\n"
         "  • Price Escalation (index-linked)\n"
         "  ✓ Payment Terms (pre-selected)\n"
         "  • Insurance & Liability\n"
         "  • Dispute Resolution (LCIA/ICC)\n"
         "  • Termination Rights\n"
         "  • Confidentiality\n\n"
         "Validation: At least 1 clause must be selected.",
         size=11, color=DARK_TEXT)

# Step 5
add_shape(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(4.5), fill_color=WHITE, border_color=GREEN, border_width=Pt(3))
add_text(slide, Inches(6.5), Inches(1.6), Inches(6.1), Inches(0.4),
         "Step 5: Optimizer Validation (Key Differentiator)", size=14, bold=True, color=GREEN)
add_text(slide, Inches(6.5), Inches(2.1), Inches(6.1), Inches(3.8),
         "The trader clicks 'Run Optimizer Validation':\n\n"
         "1. Load product group default variables (20 vars)\n"
         "2. Apply contract overrides:\n"
         "   proposed_price → nh3_price\n"
         "   proposed_freight → barge_freight\n"
         "3. Solver.solve(product_group, vars)\n"
         "   Zig/HiGHS linear program execution\n"
         "4. Display results in 4-card grid:\n"
         "   Status | Profit | Tons | ROI\n"
         "5. Show route allocations (Route 1, 2, ...)\n"
         "6. PostsolveExplainer.explain_all()\n"
         "   Claude AI generates 3-5 sentence explanation\n\n"
         "Outcomes:\n"
         "  • OPTIMAL — terms are economically viable\n"
         "  • INFEASIBLE — terms may need adjustment\n"
         "  • UNAVAILABLE — manual validation permitted\n\n"
         "The solver step is recommended but NOT required.",
         size=11, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Steps 6-7 Detail
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Steps 6-7: Review & Approval", 5)

# Step 6
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.0), fill_color=WHITE, border_color=ACCENT)
add_text(slide, Inches(0.7), Inches(1.6), Inches(5.1), Inches(0.4),
         "Step 6: Review Summary", size=14, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(2.1), Inches(5.1), Inches(3.0),
         "Consolidated view of all data:\n\n"
         "• Product Group label\n"
         "• Counterparty & Commodity\n"
         "• Delivery window dates\n"
         "• Quantity + unit, Price, Freight\n"
         "• Payment terms, Incoterm\n"
         "• Selected clauses (as tags)\n"
         "• Optimizer status, profit, ROI (if run)\n\n"
         "No edits here — trader navigates back\n"
         "to modify any step. Always valid to proceed.",
         size=11, color=DARK_TEXT)

# Step 7
add_shape(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(4.0), fill_color=WHITE, border_color=GREEN, border_width=Pt(3))
add_text(slide, Inches(6.5), Inches(1.6), Inches(6.1), Inches(0.4),
         "Step 7: Approval Submission", size=14, bold=True, color=GREEN)
add_text(slide, Inches(6.5), Inches(2.1), Inches(6.1), Inches(3.0),
         "Single Repo.transaction creates 5 records:\n\n"
         "1. CmContractNegotiation\n"
         "   Full step_data, solver snapshot, trader ID\n\n"
         "2. CmContractNegotiationEvent\n"
         "   type: submitted_for_approval\n\n"
         "3. CmContract\n"
         "   ref: CTR-{PREFIX}-{YYMMDDHHMM}-{hex}\n"
         "   status: pending_approval\n\n"
         "4. CmContractVersion (v1)\n"
         "   Full terms snapshot\n\n"
         "5. CmContractApproval\n"
         "   approver: trading_manager, status: pending\n\n"
         "Emits: contract_submitted_for_approval",
         size=11, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Event Pipeline
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Event Pipeline & Vectorization", 6)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "Every step transition emits events through the EventEmitter → Oban → pgvector pipeline, "
         "enabling semantic search over contract history.",
         size=14, color=GREY)

events = [
    ("Step 1→2", "contract_product_group_selected", "product_group, user"),
    ("Step 2→3", "contract_counterparty_set", "counterparty, commodity"),
    ("Step 3→4", "contract_commercial_terms_set", "quantity, price, freight"),
    ("Step 4→5", "contract_clauses_selected", "selected clause count"),
    ("Step 5→6", "contract_optimizer_complete", "solver status, profit, ROI"),
    ("Step 6→7", "contract_review_complete", "review confirmation"),
    ("Optimizer", "contract_optimizer_validated", "full solver results"),
    ("Submit", "contract_submitted_for_approval", "negotiation_id, contract_id, all terms"),
]

# Table header
y = Inches(2.2)
add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.4), fill_color=NAVY)
add_text(slide, Inches(0.6), y + Inches(0.03), Inches(1.5), Inches(0.35), "Trigger", size=11, bold=True, color=WHITE)
add_text(slide, Inches(2.3), y + Inches(0.03), Inches(5.0), Inches(0.35), "Event Type", size=11, bold=True, color=WHITE)
add_text(slide, Inches(7.8), y + Inches(0.03), Inches(4.8), Inches(0.35), "Key Payload", size=11, bold=True, color=WHITE)

for i, (trigger, event_type, payload) in enumerate(events):
    y = Inches(2.7 + i * 0.42)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.38), fill_color=bg, border_color=RGBColor(0xE2, 0xE8, 0xF0))
    add_text(slide, Inches(0.6), y + Inches(0.02), Inches(1.5), Inches(0.33), trigger, size=10, color=DARK_TEXT)
    add_text(slide, Inches(2.3), y + Inches(0.02), Inches(5.0), Inches(0.33), event_type, size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(7.8), y + Inches(0.02), Inches(4.8), Inches(0.33), payload, size=10, color=GREY)

# Pipeline flow
y_pipe = Inches(6.0)
pipe_items = [
    ("Contract Wizard", ACCENT), ("EventEmitter", GOLD), ("Oban Worker", ORANGE),
    ("Embeddings", PURPLE), ("pgvector", PURPLE), ("Semantic Search", GREEN)
]
for i, (label, color) in enumerate(pipe_items):
    left = Inches(0.5 + i * 2.1)
    box = add_shape(slide, left, y_pipe, Inches(1.8), Inches(0.5), fill_color=WHITE, border_color=color, border_width=Pt(2))
    add_text(slide, left + Inches(0.05), y_pipe + Inches(0.08), Inches(1.7), Inches(0.35),
             label, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Data Model
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Data Model", 7)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "Five interconnected tables created in a single transaction at approval time.",
         size=14, color=GREY)

tables = [
    ("cm_contract_negotiations", "product_group, reference_number,\ncounterparty, commodity, status,\ncurrent_step, step_data, quantity,\ndelivery_window, proposed_price,\nsolver_snapshot, trader_id", ACCENT, Inches(0.5), Inches(2.0)),
    ("cm_contract_negotiation_events", "event_type, step_number,\nactor, summary, details", ORANGE, Inches(4.8), Inches(2.0)),
    ("cm_contracts", "contract_reference, counterparty,\ncommodity, status, current_version,\nterms, selected_clause_ids,\nrequires_approval_from", ACCENT, Inches(0.5), Inches(4.2)),
    ("cm_contract_versions", "version_number, terms_snapshot,\nchange_summary, created_by", GREEN, Inches(4.8), Inches(4.2)),
    ("cm_contract_approvals", "approver_id, approval_status,\nnotes, decided_at", GOLD, Inches(9.0), Inches(4.2)),
]

for name, fields, color, left, top in tables:
    add_shape(slide, left, top, Inches(3.8), Inches(1.8), fill_color=WHITE, border_color=color, border_width=Pt(2))
    add_text(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.5), Inches(0.35),
             name, size=11, bold=True, color=color)
    add_text(slide, left + Inches(0.15), top + Inches(0.45), Inches(3.5), Inches(1.3),
             fields, size=10, color=DARK_TEXT)

# Relationship labels
add_text(slide, Inches(9.0), Inches(2.0), Inches(3.5), Inches(1.8),
         "Relationships:\n\n"
         "negotiations → has many events\n"
         "negotiations → produces contract\n"
         "contracts → has many versions\n"
         "contracts → has many approvals",
         size=11, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: Comparison vs sea.live
# ═══════════════════════════════════════════════════════════
slide = add_section_slide("Competitive Positioning vs. sea.live", 8)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "sea.live provides contract management for commodity trading. Trammo adds quantitative solver "
         "validation and AI-powered decision support.",
         size=14, color=GREY)

comparisons = [
    ("Contract Workflow", "✓ Multi-step wizard", "✓ 7-step guided wizard", "Comparable"),
    ("LP Solver Integration", "✗ Not available", "✓ HiGHS validates economics DURING contracting", "Trammo"),
    ("AI-Powered Analysis", "✗ Not available", "✓ Claude explains solver output", "Trammo"),
    ("Real-Time Data", "~ Market data feeds", "✓ 20 live variables from 10+ APIs", "Trammo"),
    ("Pre-Contract Optimization", "✗ Optimize after signing", "✓ Optimize BEFORE approval", "Trammo"),
    ("Event Vectorization", "✗ Not available", "✓ pgvector for semantic search", "Trammo"),
    ("Counterparty Management", "✓ Full KYC directory", "~ Inline entry (extensible)", "sea.live"),
    ("Solver Constraint Bridge", "✗ Contracts standalone", "✓ Approved → live constraint", "Trammo"),
]

# Table header
y = Inches(2.1)
add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.4), fill_color=NAVY)
for j, (header, left_pos) in enumerate([("Capability", 0.6), ("sea.live", 3.5), ("Trammo Trading Desk", 6.8), ("Advantage", 11.2)]):
    add_text(slide, Inches(left_pos), y + Inches(0.03), Inches(3.0), Inches(0.35),
             header, size=11, bold=True, color=WHITE)

for i, (cap, sl, tr, adv) in enumerate(comparisons):
    y = Inches(2.55 + i * 0.45)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.42), fill_color=bg, border_color=RGBColor(0xE2, 0xE8, 0xF0))
    add_text(slide, Inches(0.6), y + Inches(0.04), Inches(2.7), Inches(0.35), cap, size=10, bold=True, color=NAVY)
    sl_color = GREEN if sl.startswith("✓") else (RED if sl.startswith("✗") else ORANGE)
    add_text(slide, Inches(3.5), y + Inches(0.04), Inches(3.0), Inches(0.35), sl, size=10, color=sl_color)
    tr_color = GREEN if tr.startswith("✓") else (ORANGE if tr.startswith("~") else DARK_TEXT)
    add_text(slide, Inches(6.8), y + Inches(0.04), Inches(4.2), Inches(0.35), tr, size=10, color=tr_color)
    adv_color = GREEN if adv == "Trammo" else (RED if adv == "sea.live" else GREY)
    add_text(slide, Inches(11.2), y + Inches(0.04), Inches(1.5), Inches(0.35), adv, size=10, bold=True, color=adv_color)

# Key insight
callout = add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
                    fill_color=LIGHT_GREEN_BG, border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(1.0), Inches(6.22), Inches(11.1), Inches(0.5),
         "KEY INSIGHT: sea.live treats contracts as legal/admin workflow. Trammo inserts an LP solver BEFORE approval, "
         "transforming contract signing into a decision support moment.",
         size=11, bold=False, color=DARK_TEXT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Trammo Trading Desk  •  Contract Management Process Flows  •  February 2026  •  Confidential",
         size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
output = "/home/user/trading_desk/docs/contract_management_process_flows.pptx"
prs.save(output)
print(f"PowerPoint saved to {output}")
