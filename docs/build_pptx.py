#!/usr/bin/env python3
"""Build the 'AI Use Case Overview' PowerPoint presentation.

Agenda:
  1. Problem Statement Overview
  2. Use-Case Example
     - Set-the-Scene
     - Architectural Overview (process flow diagrams)
  3. Online Demo
  4. Questions
  5. Break-Out Session — Discover product managers' use-case scenarios

Styling matches the Trammo brand: dark navy (#0b1d3a) backgrounds, light blue
(#64B5F6) accent text, white headings, gold (#e6a817) highlights.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SCRIPT_DIR = Path(__file__).parent
IMG_DIR = SCRIPT_DIR / "diagram_images"
OUTPUT = SCRIPT_DIR / "AI_Use_Case_Overview.pptx"

# Trammo brand colours
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x0F, 0x20, 0x44)
ACCENT_BLUE = RGBColor(0x64, 0xB5, 0xF6)
MID_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
BRIGHT_BLUE = RGBColor(0x19, 0x76, 0xD2)
LIGHT_BLUE = RGBColor(0x90, 0xCA, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xC4, 0xDE)
FOOTER_GREY = RGBColor(0x60, 0x7D, 0x8B)
GOLD = RGBColor(0xE6, 0xA8, 0x17)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BODY_DARK = RGBColor(0x45, 0x5A, 0x64)
BODY_TEXT = RGBColor(0x1A, 0x23, 0x32)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=Pt(14),
                       default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT,
                       line_spacing=1.15):
    """Add a text box with multiple paragraphs. Each line is (text, size, color, bold) or just text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else default_bold

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.alignment = alignment
        p.space_after = Pt(4)
        if p.runs:
            run = p.runs[0]
            run.font.size = size
            run.font.color.rgb = color
            run.font.bold = bold

    return txBox


def add_footer(slide, left_text="Trammo  |  AI Use Case Overview",
               right_text="Confidential"):
    add_rect(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
                 left_text, Pt(9), FOOTER_GREY)
    add_text_box(slide, Inches(8.5), Inches(7.1), Inches(4.3), Inches(0.3),
                 right_text, Pt(9), FOOTER_GREY, alignment=PP_ALIGN.RIGHT)


def add_top_accent(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)


def add_section_header_slide(prs, section_num, title, subtitle=None):
    """Dark navy slide with centred section number and title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)
    add_top_accent(slide)

    add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(0.6),
                 f"SECTION {section_num}", Pt(14), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11.33), Inches(1.2),
                 title, Pt(40), WHITE, True, PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.33), Inches(0.8),
                     subtitle, Pt(16), LIGHT_BLUE, False, PP_ALIGN.CENTER)

    add_footer(slide)
    return slide


def add_diagram_image(slide, img_path, max_width=Inches(11.5), max_height=Inches(5.2),
                      center_x=True, top=Inches(1.6)):
    """Add a diagram image scaled to fit, centred horizontally."""
    img = Image.open(img_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Fit within bounds
    w = max_width
    h = int(w / aspect)
    if h > max_height:
        h = max_height
        w = int(h * aspect)

    left = (SLIDE_W - w) // 2 if center_x else Inches(0.9)
    pic = slide.shapes.add_picture(str(img_path), left, top, w, h)
    return pic


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ──────────────────────────────────────────────────────────────
    # SLIDE 1: TITLE
    # ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(11), Inches(0.5),
                 "TRAMMO TRADING TECHNOLOGY", Pt(12), ACCENT_BLUE, True)
    add_text_box(slide, Inches(0.7), Inches(1.9), Inches(11), Inches(1.2),
                 "AI Use Case Overview", Pt(44), WHITE, True)
    add_text_box(slide, Inches(0.7), Inches(3.1), Inches(9), Inches(0.5),
                 "Decision Support Platform for Commodity Trading", Pt(16), LIGHT_BLUE)
    add_text_box(slide, Inches(0.7), Inches(3.8), Inches(9), Inches(1.0),
                 "How AI-augmented optimisation transforms real-time trading decisions "
                 "across NH3 Domestic Barge, NH3 International, Sulphur, and Petcoke.",
                 Pt(11), LIGHT_GREY)

    # Product group pills
    groups = ["NH3 Domestic Barge", "NH3 International", "Sulphur International", "Petcoke"]
    x = Inches(0.7)
    for g in groups:
        shape = add_rect(slide, x, Inches(5.0), Inches(1.8), Inches(0.35), MID_BLUE, ACCENT_BLUE)
        shape.text_frame.paragraphs[0].text = g
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(9)
        run.font.color.rgb = ACCENT_BLUE
        run.font.bold = True
        x += Inches(2.0)

    add_rect(slide, Inches(0), Inches(6.85), SLIDE_W, Inches(0.65), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(11), Inches(0.3),
                 "Confidential \u2014 Internal Use Only  |  February 2026", Pt(9), FOOTER_GREY)

    # ──────────────────────────────────────────────────────────────
    # SLIDE 2: AGENDA
    # ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.9), Inches(0.5), Inches(10), Inches(0.7),
                 "Agenda", Pt(32), DARK_BLUE, True)

    agenda_items = [
        ("1", "Problem Statement Overview",
         "The challenges of manual commodity trading decisions"),
        ("2", "Use-Case Example",
         "Set-the-Scene  \u00b7  Architectural Overview"),
        ("3", "Online Demo",
         "Live walkthrough of the Trading Desk platform"),
        ("4", "Questions",
         "Open discussion"),
        ("5", "Break-Out Session",
         "Discover product managers\u2019 use-case scenarios"),
    ]

    y = Inches(1.6)
    for num, title, desc in agenda_items:
        # Number circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.5), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRIGHT_BLUE
        shape.line.fill.background()
        shape.text_frame.paragraphs[0].text = num
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True

        add_text_box(slide, Inches(1.7), y - Inches(0.02), Inches(8), Inches(0.4),
                     title, Pt(18), BRIGHT_BLUE, True)
        add_text_box(slide, Inches(1.7), y + Inches(0.38), Inches(8), Inches(0.3),
                     desc, Pt(11), BODY_DARK)
        y += Inches(1.05)

    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # SECTION 1: PROBLEM STATEMENT OVERVIEW
    # ──────────────────────────────────────────────────────────────
    add_section_header_slide(prs, "1", "Problem Statement Overview",
                             "Why commodity traders need AI-augmented decision support")

    # Problem slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.5),
                 "THE CHALLENGE", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.8), Inches(10), Inches(0.6),
                 "Manual Trading Decisions Are Slow, Fragmented & Risky",
                 Pt(24), DARK_BLUE, True)

    problems = [
        ("10+ Data Sources", "River gauges, weather, locks, market prices, freight rates, "
         "terminal inventory, fleet, capital, vessels, tides \u2014 all checked manually."),
        ("20 Variables", "Traders mentally juggle 20 interdependent variables across "
         "environmental, operational, and commercial dimensions."),
        ("Contract Complexity", "Active contracts impose price fixes, volume floors/caps, "
         "and penalty clauses that must be respected in every decision."),
        ("Decision Latency", "Assembling the data picture, checking contracts, and running "
         "spreadsheet scenarios takes hours \u2014 by which time conditions have changed."),
    ]

    y = Inches(1.7)
    for i, (title, desc) in enumerate(problems):
        # Coloured left bar
        bar_color = [BRIGHT_BLUE, GOLD, GREEN, RGBColor(0x8E, 0x44, 0xAD)][i]
        add_rect(slide, Inches(0.9), y, Inches(0.06), Inches(1.0), bar_color)
        add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(4.5), Inches(0.3),
                     title, Pt(14), DARK_BLUE, True)
        add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(4.5), Inches(0.6),
                     desc, Pt(10), BODY_DARK)
        y += Inches(1.15)

    # Solution box
    add_rect(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.6),
             RGBColor(0xEB, 0xF5, 0xFB), BRIGHT_BLUE)
    add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5.4), Inches(0.4),
                 "THE SOLUTION", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(6.8), Inches(2.4), Inches(5.4), Inches(0.5),
                 "AI-Augmented Trading Desk", Pt(18), DARK_BLUE, True)

    solutions = [
        "Continuously ingests all 10+ data sources automatically",
        "Applies contract obligations via AI clause extraction",
        "Runs LP optimisation to find profit-maximising allocation",
        "Monte Carlo (1,000 scenarios) for risk-adjusted confidence",
        "AI-generated analyst notes explain every recommendation",
        "Auto-runner triggers re-optimisation on material changes",
        "Decision latency reduced from hours to seconds",
    ]
    sol_y = Inches(3.1)
    for s in solutions:
        add_text_box(slide, Inches(7.0), sol_y, Inches(5.2), Inches(0.25),
                     f"\u2713  {s}", Pt(10), RGBColor(0x1B, 0x5E, 0x20))
        sol_y += Inches(0.3)

    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # SECTION 2: USE-CASE EXAMPLE
    # ──────────────────────────────────────────────────────────────
    add_section_header_slide(prs, "2", "Use-Case Example",
                             "Set-the-Scene  \u00b7  Architectural Overview")

    # -- 2a: Set the Scene --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "SET THE SCENE", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(11), Inches(0.6),
                 "NH3 Domestic Barge \u2014 A Typical Trading Morning",
                 Pt(24), DARK_BLUE, True)

    add_text_box(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(0.8),
                 "It\u2019s 7:00 AM. Overnight, river stage at Baton Rouge dropped 2 ft, "
                 "lock delays doubled to 12 hours, and NOLA ammonia prices jumped $8/ton. "
                 "Two active contracts have take-or-pay obligations expiring this week. "
                 "The trader needs to decide: which terminals to load, which routes to run, "
                 "and how many barges to commit \u2014 right now.",
                 Pt(13), BODY_DARK)

    # Data snapshot cards
    cards = [
        ("River Stage", "22.5 ft (\u2193 2 ft)", "USGS Baton Rouge", RGBColor(0x2E, 0x7D, 0x32)),
        ("Lock Delays", "12.0 hrs (\u2191 6 hrs)", "USACE Lower Miss.", GOLD),
        ("NOLA NH3 Price", "$348/t (\u2191 $8)", "Argus/ICIS", BRIGHT_BLUE),
        ("Meredosia Inv.", "8,200 t", "Insight TMS", RGBColor(0xE6, 0x51, 0x00)),
        ("Available Barges", "14", "Fleet TMS", RGBColor(0x8E, 0x44, 0xAD)),
        ("Working Capital", "$3.2M", "SAP S/4HANA", FOOTER_GREY),
    ]

    x = Inches(0.7)
    for title, value, source, color in cards:
        add_rect(slide, x, Inches(2.7), Inches(1.85), Inches(0.06), color)
        add_rect(slide, x, Inches(2.76), Inches(1.85), Inches(1.3),
                 RGBColor(0xFA, 0xFB, 0xFD), RGBColor(0xE2, 0xE8, 0xF0))
        add_text_box(slide, x + Inches(0.1), Inches(2.85), Inches(1.65), Inches(0.25),
                     title, Pt(9), FOOTER_GREY, True)
        add_text_box(slide, x + Inches(0.1), Inches(3.15), Inches(1.65), Inches(0.35),
                     value, Pt(16), DARK_BLUE, True)
        add_text_box(slide, x + Inches(0.1), Inches(3.6), Inches(1.65), Inches(0.2),
                     source, Pt(8), FOOTER_GREY)
        x += Inches(2.05)

    # What happens next
    add_text_box(slide, Inches(0.9), Inches(4.5), Inches(10), Inches(0.3),
                 "WHAT THE PLATFORM DOES IN SECONDS:", Pt(12), BRIGHT_BLUE, True)

    steps = [
        ("1", "Detects material changes", "Auto-runner flags river + lock + price deltas"),
        ("2", "Reads active contracts", "AI extracts take-or-pay and penalty clauses"),
        ("3", "Runs Monte Carlo", "1,000 scenarios with current conditions"),
        ("4", "Delivers recommendation", "\"Strong Go on Meredosia\u2192StL at 2,100 MT\""),
    ]
    x = Inches(0.7)
    for num, title, desc in steps:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, Inches(5.0), Inches(0.35), Inches(0.35)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRIGHT_BLUE
        shape.line.fill.background()
        shape.text_frame.paragraphs[0].text = num
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = shape.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE
        r.font.bold = True

        add_text_box(slide, x + Inches(0.5), Inches(4.93), Inches(2.3), Inches(0.3),
                     title, Pt(11), DARK_BLUE, True)
        add_text_box(slide, x + Inches(0.5), Inches(5.22), Inches(2.3), Inches(0.3),
                     desc, Pt(9), BODY_DARK)
        x += Inches(3.0)

    add_footer(slide)

    # -- 2b: Architectural Overview - System Overview --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "System Overview", Pt(24), DARK_BLUE, True)
    add_diagram_image(slide, IMG_DIR / "01_system_overview.png",
                      max_width=Inches(8), max_height=Inches(5.5), top=Inches(1.4))
    add_footer(slide)

    # -- 2b: Data Sources & Variables --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "10+ External Data Sources \u2192 20 Solver Variables", Pt(24), DARK_BLUE, True)

    # Two-column: sources list + variables table summary
    sources = [
        ("Environmental", "USGS River Gauges, NOAA Weather, USACE Lock Status", RGBColor(0x2E, 0x7D, 0x32)),
        ("Operational", "Insight TMS, Fleet TMS, SAP S/4HANA", GOLD),
        ("Commercial", "Argus/ICIS Markets, Broker Freight, EIA Nat Gas", RGBColor(0x8E, 0x44, 0xAD)),
        ("Supplementary", "AIS Vessel Tracking, NOAA Tides & Currents", FOOTER_GREY),
    ]

    y = Inches(1.5)
    for cat, items, color in sources:
        add_rect(slide, Inches(0.9), y, Inches(0.06), Inches(0.7), color)
        add_text_box(slide, Inches(1.2), y, Inches(4.5), Inches(0.3),
                     cat, Pt(13), color, True)
        add_text_box(slide, Inches(1.2), y + Inches(0.3), Inches(4.5), Inches(0.4),
                     items, Pt(10), BODY_DARK)
        y += Inches(0.85)

    # Variable groups on right side
    var_groups = [
        ("Environment (6 vars)", "river_stage, lock_hrs, temp_f, wind_mph, vis_mi, precip_in"),
        ("Operations (5 vars)", "inv_mer, inv_nio, mer_outage, nio_outage, barge_count"),
        ("Commercial (9 vars)", "nola_buy, sell_stl, sell_mem, fr_mer_stl, fr_mer_mem, "
         "fr_nio_stl, fr_nio_mem, nat_gas, working_cap"),
    ]

    add_text_box(slide, Inches(6.5), Inches(1.3), Inches(5), Inches(0.4),
                 "20 SOLVER VARIABLES", Pt(12), BRIGHT_BLUE, True)
    y = Inches(1.8)
    for group, vars_text in var_groups:
        add_rect(slide, Inches(6.5), y, Inches(5.8), Inches(0.9),
                 RGBColor(0xFA, 0xFB, 0xFD), RGBColor(0xE2, 0xE8, 0xF0))
        add_text_box(slide, Inches(6.7), y + Inches(0.05), Inches(5.4), Inches(0.3),
                     group, Pt(11), DARK_BLUE, True)
        add_text_box(slide, Inches(6.7), y + Inches(0.35), Inches(5.4), Inches(0.5),
                     vars_text, Pt(9), FOOTER_GREY)
        y += Inches(1.0)

    add_text_box(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.5),
                 "Every solve packs all 20 variables into a compact 160-byte binary "
                 "and sends it to the Zig/HiGHS solver engine.",
                 Pt(11), BRIGHT_BLUE, bold=True)
    add_footer(slide)

    # -- 2b: Data Ingestion --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Data Ingestion & Live State", Pt(24), DARK_BLUE, True)
    add_diagram_image(slide, IMG_DIR / "02_data_ingestion.png",
                      max_width=Inches(11), max_height=Inches(5.2), top=Inches(1.4))
    add_footer(slide)

    # -- 2b: Solve Pipeline --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "The 4-Phase Solve Pipeline", Pt(24), DARK_BLUE, True)
    add_diagram_image(slide, IMG_DIR / "03_solve_pipeline.png",
                      max_width=Inches(5), max_height=Inches(5.5), top=Inches(1.4))
    add_footer(slide)

    # -- 2b: Pre-Solve LLM --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Pre-Solve: LLM Contract Framing", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4),
                 "AI reads active contracts and adjusts solver variables to respect "
                 "price fixes, take-or-pay volumes, and penalty clauses.",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "04_presolve_llm.png",
                      max_width=Inches(11), max_height=Inches(4.5), top=Inches(1.9))
    add_footer(slide)

    # -- 2b: Solver Engine --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "The Solver Engine \u2014 LP + Monte Carlo", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4),
                 "High-performance Zig/HiGHS linear-programming solver. Single solve for optimal "
                 "allocation; Monte Carlo (1,000 scenarios) for risk-adjusted confidence signals.",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "05_solver_engine.png",
                      max_width=Inches(11), max_height=Inches(4.5), top=Inches(1.9))
    add_footer(slide)

    # -- 2b: Post-Solve LLM --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Post-Solve: AI Analyst Explanation", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.7),
                 "After the solver returns, AI writes a concise analyst note:\n"
                 "\u201cThe solver favors Meredosia-to-StL at 2,100 MT, capturing a $40/t margin. "
                 "Memphis allocation limited by barge availability (shadow price: $12.50/barge).\u201d",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "06_postsolve_llm.png",
                      max_width=Inches(11), max_height=Inches(4.0), top=Inches(2.2))
    add_footer(slide)

    # -- 2b: Auto-Runner --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Auto-Runner: Continuous Monitoring", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4),
                 "When live data changes materially, the auto-runner triggers Monte Carlo, "
                 "notifies traders, and records the result on an immutable audit chain.",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "07_auto_runner.png",
                      max_width=Inches(5), max_height=Inches(5.0), top=Inches(1.8))
    add_footer(slide)

    # -- 2b: Contract Lifecycle --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Contract Lifecycle", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4),
                 "Upload \u2192 AI clause extraction \u2192 Legal review \u2192 SAP validation "
                 "\u2192 Trading gate \u2192 Solver constraints.",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "08_contract_lifecycle.png",
                      max_width=Inches(11), max_height=Inches(4.5), top=Inches(1.9))
    add_footer(slide)

    # -- 2b: Decision Ledger --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Trader Decision Ledger", Pt(24), DARK_BLUE, True)
    add_diagram_image(slide, IMG_DIR / "09_decision_ledger.png",
                      max_width=Inches(6), max_height=Inches(5.2), top=Inches(1.4))
    add_footer(slide)

    # -- 2b: What-If Analysis --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "What-If Analysis Workflow", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4),
                 "5-step wizard: Variables \u2192 Contracts \u2192 Natural-language prompt "
                 "\u2192 AI framing \u2192 Solver results with explanation.",
                 Pt(11), BODY_DARK)
    add_diagram_image(slide, IMG_DIR / "10_whatif_analysis.png",
                      max_width=Inches(11), max_height=Inches(4.5), top=Inches(1.9))
    add_footer(slide)

    # -- 2b: End-to-End Flow --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "ARCHITECTURAL OVERVIEW", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "End-to-End: From Raw Data to Trading Decision",
                 Pt(24), DARK_BLUE, True)
    add_diagram_image(slide, IMG_DIR / "11_end_to_end.png",
                      max_width=Inches(8), max_height=Inches(5.5), top=Inches(1.4))
    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # SECTION 3: ONLINE DEMO
    # ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_top_accent(slide)

    add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(0.6),
                 "SECTION 3", Pt(14), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11.33), Inches(1.2),
                 "Online Demo", Pt(40), WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.33), Inches(0.8),
                 "Live walkthrough of the Trading Desk platform",
                 Pt(16), LIGHT_BLUE, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9.33), Inches(0.5),
                 "http://localhost:4111", Pt(14), GOLD, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # -- Demo Prompts: What-If Scenarios --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_accent(slide)
    add_text_box(slide, Inches(0.9), Inches(0.3), Inches(10), Inches(0.4),
                 "DEMO \u2014 WHAT-IF SCENARIOS", Pt(12), BRIGHT_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.5),
                 "Copy & Paste Prompts for the Live Demo", Pt(24), DARK_BLUE, True)
    add_text_box(slide, Inches(0.9), Inches(1.25), Inches(10), Inches(0.3),
                 "Enter these into the What-If prompt box to demonstrate AI-framed scenario analysis:",
                 Pt(11), BODY_DARK)

    demo_prompts = [
        ("\u2460  River Drop + Lock Surge",
         "What if river stage drops to 7 feet and lock delays increase to 24 hours? "
         "Which contracts can still operate and what happens to the Koch demurrage exposure?"),
        ("\u2461  NOLA Price Spike",
         "What if NOLA ammonia prices jump to $380/MT? Does the CF Industries $365 cap "
         "still hold and can we meet the Mosaic $415 floor at St. Louis?"),
        ("\u2462  Barge Shortage",
         "What if we lose 3 barges and only have 2 available? Can we still meet the "
         "Simplot quarterly minimum of 6,750 MT across both destinations?"),
        ("\u2463  Terminal Outage",
         "What if Meredosia terminal goes into unplanned outage? How do we reroute "
         "the Mosaic St. Louis and Simplot deliveries from Geismar?"),
        ("\u2464  Working Capital Squeeze",
         "What if working capital drops to $2.5M? The CF Industries LC alone requires "
         "$4M \u2014 which contracts do we prioritize and what\u2019s the profit impact?"),
        ("\u2465  Combined Stress Test",
         "What if lock delays double to 24 hours, we lose 2 barges, and NOLA price "
         "rises to $360? Run a Monte Carlo \u2014 is this still a go across the portfolio?"),
    ]

    y = Inches(1.7)
    for label, prompt in demo_prompts:
        add_rect(slide, Inches(0.7), y, Inches(11.7), Inches(0.82),
                 RGBColor(0xFA, 0xFB, 0xFD), RGBColor(0xE2, 0xE8, 0xF0))
        add_text_box(slide, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.25),
                     label, Pt(10), BRIGHT_BLUE, True)
        add_text_box(slide, Inches(0.9), y + Inches(0.3), Inches(11.2), Inches(0.5),
                     prompt, Pt(10), BODY_TEXT)
        y += Inches(0.9)

    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # SECTION 4: QUESTIONS
    # ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_top_accent(slide)

    add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(0.6),
                 "SECTION 4", Pt(14), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11.33), Inches(1.2),
                 "Questions", Pt(40), WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.33), Inches(0.8),
                 "Open discussion", Pt(16), LIGHT_BLUE, False, PP_ALIGN.CENTER)
    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # SECTION 5: BREAK-OUT SESSION
    # ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_top_accent(slide)

    add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.6),
                 "SECTION 5", Pt(14), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.33), Inches(1.2),
                 "Break-Out Session", Pt(40), WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10.33), Inches(0.8),
                 "Discover product managers\u2019 use-case scenarios",
                 Pt(16), LIGHT_BLUE, False, PP_ALIGN.CENTER)

    prompts = [
        "What trading decisions does your team make daily?",
        "Which data sources do you currently check manually?",
        "Where are the biggest bottlenecks in your workflow?",
        "What would \u201creal-time optimisation\u201d look like for your product?",
    ]
    y = Inches(4.5)
    for prompt in prompts:
        add_text_box(slide, Inches(2.5), y, Inches(8.33), Inches(0.35),
                     f"\u2022  {prompt}", Pt(14), LIGHT_BLUE)
        y += Inches(0.45)

    add_footer(slide)

    # ──────────────────────────────────────────────────────────────
    # Save
    # ──────────────────────────────────────────────────────────────
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
