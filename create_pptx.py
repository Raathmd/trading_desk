"""
Generate executive_overview.pptx from the Trading Desk Executive Overview HTML content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color palette (matches the HTML) ─────────────────────────────
NAVY       = RGBColor(0x0f, 0x20, 0x44)
BLUE       = RGBColor(0x19, 0x76, 0xd2)
LIGHT_BLUE = RGBColor(0xe3, 0xf2, 0xfd)
DARK_BLUE  = RGBColor(0x0d, 0x47, 0xa1)
TEXT_DARK  = RGBColor(0x1a, 0x23, 0x32)
TEXT_MID   = RGBColor(0x45, 0x5a, 0x64)
TEXT_LIGHT = RGBColor(0x60, 0x7d, 0x8b)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GREEN_DARK = RGBColor(0x2e, 0x7d, 0x32)
GREEN_LIGHT= RGBColor(0xe8, 0xf5, 0xe9)
AMBER_DARK = RGBColor(0xe6, 0x51, 0x00)
AMBER_LIGHT= RGBColor(0xff, 0xf3, 0xe0)
PURPLE_DARK= RGBColor(0x6a, 0x1b, 0x9a)
PURPLE_LIGHT=RGBColor(0xf3, 0xe5, 0xf5)
BABY_BLUE  = RGBColor(0x64, 0xb5, 0xf6)
STEEL_BLUE = RGBColor(0x90, 0xca, 0xf9)
LIGHT_GREY = RGBColor(0xf5, 0xf7, 0xfa)
BORDER_GREY= RGBColor(0xb0, 0xbe, 0xc5)

# ── Presentation setup ────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout

def add_slide():
    return prs.slides.add_slide(BLANK)

def tf(shape):
    return shape.text_frame

def add_rect(slide, left, top, width, height,
             fill=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="",
                font_size=Pt(11), bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.text_frame.word_wrap = wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para(tf_obj, text, font_size=Pt(11), bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False,
             space_before=Pt(4)):
    """Append a paragraph to an existing text frame."""
    from pptx.util import Pt as ptConv
    para = tf_obj.add_paragraph()
    para.alignment = align
    para.space_before = space_before
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return para

def section_header(slide, number, title, intro, label_color=BLUE):
    """Draw the standard section header band."""
    # Top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, fill=BLUE)
    # Label
    add_textbox(slide, 0.4, 0.15, 3, 0.3,
                f"SECTION {number}", font_size=Pt(9),
                bold=True, color=label_color)
    # Title
    add_textbox(slide, 0.4, 0.4, 12.5, 0.7, title,
                font_size=Pt(26), bold=True, color=NAVY)
    # Intro text box with blue bottom border effect — just text
    add_textbox(slide, 0.4, 1.1, 12.3, 0.8, intro,
                font_size=Pt(11), color=TEXT_MID, wrap=True)
    # Thin blue rule under intro
    add_rect(slide, 0.4, 1.85, 9, 0.03, fill=BLUE)

def footer(slide, section_text):
    add_rect(slide, 0, 7.2, 13.33, 0.02, fill=BORDER_GREY)
    add_textbox(slide, 0.4, 7.25, 6, 0.2,
                "Trammo Trading Desk — Executive Overview",
                font_size=Pt(8), color=TEXT_LIGHT)
    add_textbox(slide, 7, 7.25, 6, 0.2, section_text,
                font_size=Pt(8), color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

def callout_box(slide, left, top, width, height, text,
                border_color=BLUE, bg=LIGHT_BLUE, text_color=DARK_BLUE):
    box = add_rect(slide, left, top, width, height,
                   fill=bg, line_color=border_color, line_width=Pt(2))
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(height - 0.2)
    )
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = text_color

def small_box(slide, left, top, width, height, title, body,
              border=BORDER_GREY, bg=WHITE, title_color=DARK_BLUE):
    add_rect(slide, left, top, width, height,
             fill=bg, line_color=border, line_width=Pt(1.5))
    # title
    add_textbox(slide, left + 0.1, top + 0.08, width - 0.2, 0.28,
                title, font_size=Pt(10), bold=True, color=title_color)
    # body
    tb = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.35),
        Inches(width - 0.2), Inches(height - 0.45)
    )
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = body
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_MID

def flow_box(slide, left, top, width, height, title, subtitle="",
             border=BLUE, bg=LIGHT_BLUE, fg=DARK_BLUE):
    add_rect(slide, left, top, width, height,
             fill=bg, line_color=border, line_width=Pt(2))
    add_textbox(slide, left + 0.05, top + 0.07, width - 0.1,
                height * 0.5, title,
                font_size=Pt(9), bold=True, color=fg,
                align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, left + 0.05, top + height * 0.5,
                    width - 0.1, height * 0.5, subtitle,
                    font_size=Pt(7.5), color=fg,
                    align=PP_ALIGN.CENTER)

def arrow(slide, left, top, width=0.3, horizontal=True):
    """Draw a simple text arrow."""
    symbol = "→" if horizontal else "↓"
    add_textbox(slide, left, top, width, 0.3, symbol,
                font_size=Pt(16), color=TEXT_LIGHT,
                align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
slide = add_slide()

# Full dark background
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)

# Label
add_textbox(slide, 0.7, 0.7, 8, 0.4, "TRAMMO TRADING TECHNOLOGY",
            font_size=Pt(9), bold=True, color=BABY_BLUE)

# Main title
tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(9), Inches(2))
tf_obj = tb.text_frame
tf_obj.word_wrap = True
p = tf_obj.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Trading Desk"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = WHITE

p2 = tf_obj.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
run2 = p2.add_run()
run2.text = "Executive Overview"
run2.font.size = Pt(44)
run2.font.bold = True
run2.font.color.rgb = WHITE

# Subtitle
add_textbox(slide, 0.7, 3.15, 9, 0.45,
            "Decision Support Platform for Commodity Trading",
            font_size=Pt(16), color=STEEL_BLUE)

# Description
add_textbox(slide, 0.7, 3.7, 9, 1.0,
            "This document explains how the Trammo Trading Desk platform integrates "
            "real-time market data, contractual obligations, and supply chain constraints "
            "to identify and record optimal trading decisions — and how it keeps traders "
            "informed when conditions shift.",
            font_size=Pt(11), color=RGBColor(0xb0, 0xc4, 0xde), wrap=True)

# Product tags
tags = ["NH3 Domestic Barge", "NH3 International", "Sulphur International", "Petcoke"]
tag_x = 0.7
for tag in tags:
    w = len(tag) * 0.095 + 0.35
    add_rect(slide, tag_x, 4.85, w, 0.35,
             fill=RGBColor(0x1a, 0x3a, 0x6b),
             line_color=RGBColor(0x2e, 0x5b, 0xa0), line_width=Pt(1))
    add_textbox(slide, tag_x + 0.1, 4.88, w - 0.2, 0.28, tag,
                font_size=Pt(9), bold=True, color=BABY_BLUE)
    tag_x += w + 0.15

# Thin divider
add_rect(slide, 0.7, 5.45, 11.5, 0.02, fill=RGBColor(0x1e, 0x3a, 0x5f))

# Footer meta
add_textbox(slide, 0.7, 5.55, 5.5, 0.3,
            "Confidential — Internal Use Only",
            font_size=Pt(9), color=TEXT_LIGHT)
add_textbox(slide, 7.5, 5.55, 5, 0.3, "February 2026",
            font_size=Pt(9), color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, 13.33, 0.08, fill=BLUE)

add_textbox(slide, 0.5, 0.2, 10, 0.6, "Contents",
            font_size=Pt(22), bold=True, color=NAVY)

toc_items = [
    ("1", "What the Platform Does",
     "How it supports each product group"),
    ("2", "Data Ingestion & the Live Model",
     "Nine live data sources, polling schedule, variable assembly"),
    ("3", "How the Operational Model is Formed",
     "From raw data to a mathematical representation of the business"),
    ("4", "Contract Ingestion & SAP Integration",
     "Contractual constraints, SAP validation, open-position sync"),
    ("5", "How the Solver Finds the Optimal Solution",
     "Linear programming, Monte Carlo, shadow prices, decision signals"),
    ("6", "The Auto Trader",
     "Continuous monitoring, delta detection, automatic re-optimisation"),
    ("7", "Trader Notifications & Threshold Alerts",
     "Email, Slack, Teams — what triggers an alert and what it contains"),
]

y = 1.0
for num, title, sub in toc_items:
    # Number badge
    add_rect(slide, 0.5, y, 0.45, 0.45, fill=NAVY)
    add_textbox(slide, 0.5, y + 0.05, 0.45, 0.35, num,
                font_size=Pt(13), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, 1.1, y, 9, 0.3, title,
                font_size=Pt(12), bold=True, color=DARK_BLUE)
    # Sub
    add_textbox(slide, 1.1, y + 0.3, 9, 0.2, sub,
                font_size=Pt(9.5), color=TEXT_MID)
    # Rule
    add_rect(slide, 1.1, y + 0.5, 11.5, 0.01, fill=BORDER_GREY)
    y += 0.65

footer(slide, "Table of Contents")

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — SECTION 1: WHAT THE PLATFORM DOES
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "1", "What the Platform Does",
    "The Trading Desk is a real-time decision-support system that continuously ingests live market, "
    "operational, and environmental data to calculate the optimal trade allocation across every "
    "active route — and tells traders what to do, why, and when conditions change.")

# Core capability callout
callout_box(slide, 0.4, 2.05, 12.4, 1.0,
    "At any moment, the platform holds a precise mathematical picture of the current operating "
    "environment — river depth, freight rates, inventory, commodity prices, barge availability, "
    "weather, lock delays, and working capital — all drawn from live APIs. It applies every active "
    "contractual obligation and solves a linear programme to find the allocation that maximises profit.")

# Product group table header
add_textbox(slide, 0.4, 3.15, 12, 0.3, "HOW EACH PRODUCT GROUP BENEFITS",
            font_size=Pt(9), bold=True, color=BLUE)

# Table header row
add_rect(slide, 0.4, 3.45, 12.4, 0.32, fill=NAVY)
for col_x, col_w, col_text in [
    (0.45, 2.4, "Product Group"),
    (2.9, 5.0, "Key Decision Supported"),
    (7.95, 4.75, "Critical Variables"),
]:
    add_textbox(slide, col_x, 3.48, col_w, 0.26, col_text,
                font_size=Pt(9), bold=True, color=WHITE)

rows = [
    ("NH3 Domestic Barge",
     "Which terminals to load, which destinations to serve, how many barges to commit",
     "River stage, barge count, NOLA spot price, freight rates, lock delays"),
    ("NH3 International",
     "Vessel scheduling, NOLA vs export arbitrage, vessel routing optimisation",
     "Vessel positions, ocean freight, FOB/CIF spread, terminal inventory"),
    ("Sulphur International",
     "Export volume allocation, vessel nomination, penalty exposure management",
     "Refinery run rates, vessel capacity, contract penalty clauses, spot premiums"),
    ("Petcoke",
     "Destination market allocation, freight arbitrage, inventory drawdown optimisation",
     "Refinery availability, port draft, ocean freight, quality differentials"),
]
row_y = 3.77
for i, (prod, decision, crit) in enumerate(rows):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    add_rect(slide, 0.4, row_y, 12.4, 0.42, fill=bg, line_color=None)
    add_textbox(slide, 0.45, row_y + 0.04, 2.35, 0.35, prod,
                font_size=Pt(9), bold=True, color=DARK_BLUE)
    add_textbox(slide, 2.9, row_y + 0.04, 4.9, 0.35, decision,
                font_size=Pt(8.5), color=TEXT_DARK, wrap=True)
    add_textbox(slide, 7.95, row_y + 0.04, 4.7, 0.35, crit,
                font_size=Pt(8.5), color=TEXT_DARK, wrap=True)
    row_y += 0.42

# Key principle callout
callout_box(slide, 0.4, 5.6, 12.4, 0.55,
    "Key principle: Every solve incorporates the current contractual position from SAP and the live "
    "market environment simultaneously — so the recommended trade always reflects both what the "
    "market offers and what the contracts require.",
    border_color=GREEN_DARK, bg=GREEN_LIGHT, text_color=RGBColor(0x1b, 0x5e, 0x20))

# Two modes
small_box(slide, 0.4, 6.22, 6.0, 0.82,
    "Trader Mode — Scenario Exploration",
    "Open the dashboard, see 20 live variables, adjust with sliders, choose an objective, "
    "click Solve. Get route allocations, margins, shadow prices, and plain-English explanation in seconds.")
small_box(slide, 6.75, 6.22, 6.05, 0.82,
    "Agent Mode — Continuous Optimisation",
    "Auto-runner monitors 20 variables continuously. When any moves beyond its threshold, "
    "it re-solves with 1,000 Monte Carlo scenarios, records the result, commits to audit trail, and notifies the trader.")

footer(slide, "Section 1: What the Platform Does")

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — SECTION 2: DATA INGESTION
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "2", "Data Ingestion & the Live Model",
    "Nine external data sources feed the platform continuously. Each source is polled on a schedule, "
    "validated on arrival, and written into a live state store that the solver reads from at the moment it is called.")

# Source cards (2-row grid, 4 cols each)
sources = [
    ("USGS Gauges", "River stage (ft), flow (cfs) — 4 gauges from Cairo IL to Baton Rouge LA", "Every 15 min"),
    ("NOAA Weather", "Temp, wind, visibility, 3-day precipitation — 4 stations along route", "Every 30 min"),
    ("USACE Lock Status", "Cumulative lock delay hours on Lower Mississippi", "Every 30 min"),
    ("EIA — Henry Hub", "Natural gas price ($/MMBtu) — affects NH3 production economics", "Every 60 min"),
    ("Market Feeds", "NOLA buy price, St. Louis delivered, Memphis delivered ($/ton)", "Every 30 min"),
    ("Freight Broker", "4 route rates: Don→StL, Don→Mem, Geis→StL, Geis→Mem ($/ton)", "Every 60 min"),
    ("Internal Systems", "Terminal inventory, barge count, dock outages, working capital", "Every 5 min"),
    ("Vessel Tracking", "Fleet GPS, ETA, weather at vessel position, carrying stock", "Every 10 min"),
]

card_w = 2.9
card_h = 0.88
gap = 0.18
start_x = 0.4
start_y = 2.05

for idx, (name, data, freq) in enumerate(sources):
    col = idx % 4
    row = idx // 4
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + 0.1)
    add_rect(slide, x, y, card_w, card_h, fill=LIGHT_GREY, line_color=BORDER_GREY, line_width=Pt(1.5))
    add_textbox(slide, x + 0.1, y + 0.07, card_w - 0.2, 0.26, name,
                font_size=Pt(9.5), bold=True, color=TEXT_DARK)
    add_textbox(slide, x + 0.1, y + 0.32, card_w - 0.2, 0.4, data,
                font_size=Pt(8), color=TEXT_LIGHT, wrap=True)
    add_textbox(slide, x + 0.1, y + 0.68, card_w - 0.2, 0.18, freq,
                font_size=Pt(7.5), color=RGBColor(0x90, 0xa4, 0xae), bold=True)

# Arrow down to live state store
add_textbox(slide, 6.1, 4.05, 1, 0.4, "↓↓↓↓↓↓↓↓",
            font_size=Pt(14), color=BLUE, align=PP_ALIGN.CENTER)

# Live state store
flow_box(slide, 3.5, 4.45, 6.3, 0.55,
         "Live State Store (ETS in-memory)",
         "All 20 variables — always current, sub-second reads",
         border=NAVY, bg=NAVY, fg=STEEL_BLUE)

# Arrow down
add_textbox(slide, 6.4, 5.02, 0.5, 0.3, "↓", font_size=Pt(16), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Three boxes in a row
flow_box(slide, 0.4, 5.33, 3.8, 0.6, "Variable Assembly",
         "20 variables packed to binary for solver")
flow_box(slide, 4.4, 5.33, 4.2, 0.6, "Contract Constraint Bridge",
         "Active clauses tighten variable bounds",
         border=PURPLE_DARK, bg=PURPLE_LIGHT, fg=PURPLE_DARK)
flow_box(slide, 8.8, 5.33, 4.1, 0.6, "SAP Open Positions",
         "Current undelivered qty per contract",
         border=AMBER_DARK, bg=AMBER_LIGHT, fg=AMBER_DARK)

# Arrow down
add_textbox(slide, 6.4, 5.95, 0.5, 0.3, "↓", font_size=Pt(16), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Solver box
flow_box(slide, 3.5, 6.28, 6.3, 0.55,
         "Zig LP Solver (HiGHS engine)",
         "Linear programme + Monte Carlo 1,000 scenarios",
         border=BLUE, bg=NAVY, fg=STEEL_BLUE)

footer(slide, "Section 2: Data Ingestion & the Live Model")

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — SECTION 2 cont'd: THE 20 VARIABLES
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, 13.33, 0.08, fill=BLUE)
add_textbox(slide, 0.4, 0.15, 5, 0.4, "SECTION 2 (continued)",
            font_size=Pt(9), bold=True, color=BLUE)
add_textbox(slide, 0.4, 0.45, 12, 0.55, "The 20 Live Variables and Their Sources",
            font_size=Pt(20), bold=True, color=NAVY)

# Environment group
add_rect(slide, 0.4, 1.2, 12.4, 1.1, fill=LIGHT_BLUE, line_color=BLUE, line_width=Pt(1.5))
add_textbox(slide, 0.55, 1.28, 8, 0.3,
            "Environment (6) — drawn from USGS, NOAA, USACE",
            font_size=Pt(10), bold=True, color=DARK_BLUE)
env_chips = ["river_stage (ft)", "lock_hrs (delay hrs)", "temp_f (°F)", "wind_mph", "vis_mi (visibility)", "precip_in (3-day rain)"]
cx = 0.55
cy = 1.65
for chip in env_chips:
    cw = len(chip) * 0.09 + 0.25
    add_rect(slide, cx, cy, cw, 0.3, fill=LIGHT_BLUE, line_color=BLUE, line_width=Pt(1))
    add_textbox(slide, cx + 0.08, cy + 0.04, cw - 0.1, 0.22, chip,
                font_size=Pt(8.5), color=DARK_BLUE)
    cx += cw + 0.12

# Operations group
add_rect(slide, 0.4, 2.5, 12.4, 1.1, fill=GREEN_LIGHT, line_color=GREEN_DARK, line_width=Pt(1.5))
add_textbox(slide, 0.55, 2.58, 9, 0.3,
            "Operations (5) — drawn from internal ERP, fleet tracking",
            font_size=Pt(10), bold=True, color=GREEN_DARK)
ops_chips = ["inv_don (tons, Donaldsonville)", "inv_geis (tons, Geismar)", "barge_count (available fleet)", "stl_outage (dock status)", "mem_outage (dock status)"]
cx = 0.55
cy = 2.95
for chip in ops_chips:
    cw = len(chip) * 0.09 + 0.25
    add_rect(slide, cx, cy, cw, 0.3, fill=GREEN_LIGHT, line_color=GREEN_DARK, line_width=Pt(1))
    add_textbox(slide, cx + 0.08, cy + 0.04, cw - 0.1, 0.22, chip,
                font_size=Pt(8.5), color=GREEN_DARK)
    cx += cw + 0.12

# Commercial group
add_rect(slide, 0.4, 3.8, 12.4, 1.2, fill=AMBER_LIGHT, line_color=AMBER_DARK, line_width=Pt(1.5))
add_textbox(slide, 0.55, 3.88, 10, 0.3,
            "Commercial (9) — drawn from market feeds, freight brokers, EIA",
            font_size=Pt(10), bold=True, color=AMBER_DARK)
com_chips = ["nola_buy ($/t)", "sell_stl ($/t)", "sell_mem ($/t)", "fr_don_stl ($/t)", "fr_don_mem ($/t)", "fr_geis_stl ($/t)", "fr_geis_mem ($/t)", "nat_gas ($/MMBtu)", "working_cap ($)"]
cx = 0.55
cy = 4.25
for chip in com_chips:
    cw = len(chip) * 0.09 + 0.25
    if cx + cw > 12.6:
        cx = 0.55
        cy += 0.38
    add_rect(slide, cx, cy, cw, 0.3, fill=AMBER_LIGHT, line_color=AMBER_DARK, line_width=Pt(1))
    add_textbox(slide, cx + 0.08, cy + 0.04, cw - 0.1, 0.22, chip,
                font_size=Pt(8.5), color=AMBER_DARK)
    cx += cw + 0.12

# Delta detection callout
callout_box(slide, 0.4, 5.3, 12.4, 0.75,
    "Delta Detection: Each variable carries a configured sensitivity threshold "
    "(e.g., river_stage ±0.5 ft, nola_buy ±$2/ton). The auto-runner only re-solves when at "
    "least one variable moves beyond its threshold — avoiding unnecessary computation while "
    "still reacting promptly to material changes.")

# Solver outputs
add_textbox(slide, 0.4, 6.15, 12, 0.28, "SOLVER OUTPUTS",
            font_size=Pt(9), bold=True, color=BLUE)
output_items = [
    ("Route Allocations", "Tons per route,\nbarges required"),
    ("Profit & Margins", "Total P&L, $/ton\nper route"),
    ("Risk Distribution", "P5/P50/P95,\nVaR, signal"),
    ("Shadow Prices", "Value of relaxing\neach constraint"),
]
ox = 0.4
for title, sub in output_items:
    flow_box(slide, ox, 6.42, 2.95, 0.72, title, sub,
             border=GREEN_DARK, bg=GREEN_LIGHT, fg=GREEN_DARK)
    ox += 3.15

footer(slide, "Section 2: Data Ingestion & the Live Model")

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — SECTION 3: HOW THE MODEL IS FORMED
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "3", "How the Operational Model Is Formed",
    "The 20 live variables are assembled into a structured mathematical model — a linear programme — "
    "that represents the real physical and commercial constraints of the business.")

# Model components table
add_textbox(slide, 0.4, 2.0, 12, 0.3, "WHAT THE MODEL ENCODES (NH3 Domestic Barge example)",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 0.4, 2.3, 12.4, 0.32, fill=NAVY)
for col_x, col_w, col_text in [(0.45, 2.3, "Model Component"), (2.8, 4.5, "What It Represents"), (7.35, 5.35, "How Variables Shape It")]:
    add_textbox(slide, col_x, 2.33, col_w, 0.26, col_text, font_size=Pt(9), bold=True, color=WHITE)

model_rows = [
    ("Routes (4)", "Don→StL, Don→Mem, Geis→StL, Geis→Mem — each a physical path from terminal to customer",
     "Freight rates, sell prices, and buy price determine the margin per ton on each route"),
    ("Barge Capacity", "Effective tons per barge — reduced when river stage is low",
     "river_stage <12 ft → 75% capacity; 12–18 ft → 90%; ≥18 ft → 100% (1,500 MT base)"),
    ("Fleet Constraint", "Total barges committed cannot exceed available fleet",
     "barge_count sets the hard upper bound; effective capacity scales with river stage"),
    ("Inventory Constraints", "Cannot load more from a terminal than it holds",
     "inv_don and inv_geis are the supply caps for Donaldsonville and Geismar routes"),
    ("Capital Constraint", "Total product purchased cannot exceed available working capital",
     "working_cap sets the liquidity ceiling; nola_buy determines cost per ton purchased"),
    ("Outage Flags", "If a destination dock is out of service, that route's allocation is forced to zero",
     "stl_outage / mem_outage set the upper bound on relevant route variables to zero"),
    ("Nat Gas Cost", "Natural gas is a significant input cost at ammonia production terminals",
     "nat_gas adjusts the effective buy-side cost per ton, compressing or expanding margin"),
]
row_y = 2.62
for i, (comp, what, how) in enumerate(model_rows):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    h = 0.4
    add_rect(slide, 0.4, row_y, 12.4, h, fill=bg)
    add_textbox(slide, 0.45, row_y + 0.04, 2.25, h - 0.08, comp,
                font_size=Pt(8.5), bold=True, color=DARK_BLUE)
    add_textbox(slide, 2.8, row_y + 0.04, 4.4, h - 0.08, what,
                font_size=Pt(8), color=TEXT_DARK, wrap=True)
    add_textbox(slide, 7.35, row_y + 0.04, 5.3, h - 0.08, how,
                font_size=Pt(8), color=TEXT_DARK, wrap=True)
    row_y += h

# Example amber callout
callout_box(slide, 0.4, 5.5, 12.4, 0.65,
    "Example: River stage drops from 18 ft to 10 ft overnight. Effective barge capacity falls from "
    "1,500 MT to 1,125 MT. The optimal solution may shift — fewer tons to distant destinations. "
    "The auto-runner detects the drop, triggers a re-solve, and notifies the trader with the updated "
    "allocation and new P&L figure.",
    border_color=AMBER_DARK, bg=AMBER_LIGHT, text_color=RGBColor(0xbf, 0x36, 0x0c))

# Objective modes
add_textbox(slide, 0.4, 6.22, 12, 0.25, "OBJECTIVE MODES",
            font_size=Pt(9), bold=True, color=BLUE)
modes = [
    ("Max Profit", "Maximise total gross profit across all routes. Standard mode for most commercial decisions."),
    ("Min Cost", "Minimise total capital deployed while meeting delivery obligations. Used under capital constraints."),
    ("Max ROI", "Maximise return on capital deployed — useful when competing with other product group opportunities."),
    ("CVaR-Adjusted", "Maximise expected profit while penalising tail risk. Monte Carlo distribution informs this."),
]
mx = 0.4
for title, desc in modes:
    small_box(slide, mx, 6.48, 3.05, 0.82, title, desc)
    mx += 3.23

footer(slide, "Section 3: How the Operational Model Is Formed")

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — SECTION 4: CONTRACT INGESTION & SAP
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "4", "Contract Ingestion & SAP Integration",
    "Before the solver can run, the platform must know what the business is contractually obligated "
    "to deliver — and what SAP records as the current open position. Contracts are the source of "
    "truth for constraints; SAP is the source of truth for positions.")

# 5-step pipeline
steps = [
    ("1", NAVY, "Document Detection",
     "Platform computes SHA-256 hash of each contract file in SharePoint. When hash changes (new version), ingestion is triggered automatically."),
    ("2", BLUE, "AI Clause Extraction",
     "LLM pipeline extracts structured clause objects: quantity tolerances, delivery windows, penalty rates, price terms, incoterms, force majeure."),
    ("3", PURPLE_DARK, "SAP Cross-Validation",
     "Platform fetches matching contract record from SAP (OData API) and compares extracted clause values. Discrepancies flagged for operations review."),
    ("4", AMBER_DARK, "Legal Review & Approval",
     "Legal team reviews extracted clauses in the Contracts dashboard and approves or rejects. Only approved, SAP-validated contracts become active."),
    ("5", GREEN_DARK, "Activation as Solver Constraints",
     "Active clauses flow through the Constraint Bridge into the solver. Quantity minimums become floor constraints. Penalty clauses reduce effective margin."),
]
y = 2.05
for num, num_color, title, desc in steps:
    add_rect(slide, 0.4, y, 0.4, 0.4, fill=num_color)
    add_textbox(slide, 0.4, y + 0.07, 0.4, 0.28, num,
                font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, y + 0.02, 3.5, 0.26, title,
                font_size=Pt(10), bold=True, color=TEXT_DARK)
    add_textbox(slide, 0.9, y + 0.27, 5.7, 0.22, desc,
                font_size=Pt(8.5), color=TEXT_MID, wrap=True)
    if int(num) < 5:
        add_textbox(slide, 0.52, y + 0.4, 0.2, 0.15, "↓",
                    font_size=Pt(10), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    y += 0.62

# SAP fields table on right side
add_textbox(slide, 7.1, 2.0, 6, 0.3, "SAP FIELDS — SYSTEM-OF-RECORD POSITION",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 7.1, 2.3, 5.9, 0.32, fill=NAVY)
for col_x, col_w, col_text in [(7.15, 2.0, "SAP Field"), (9.2, 1.9, "Meaning"), (11.15, 1.8, "Solver Use")]:
    add_textbox(slide, col_x, 2.33, col_w, 0.26, col_text, font_size=Pt(9), bold=True, color=WHITE)

sap_rows = [
    ("ContractQuantity", "Total agreed quantity", "Maximum volume ceiling"),
    ("OpenQuantity", "Remaining undelivered qty", "Minimum allocation floor"),
    ("DeliveredQuantity", "Confirmed deliveries to date", "Verifies open position consistency"),
    ("Incoterms", "FOB, CIF, DAP, etc.", "Determines cost elements — affects margin"),
]
ry = 2.62
for i, (field, meaning, use) in enumerate(sap_rows):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    add_rect(slide, 7.1, ry, 5.9, 0.42, fill=bg)
    add_textbox(slide, 7.15, ry + 0.05, 1.95, 0.32, field,
                font_size=Pt(8.5), bold=True, color=DARK_BLUE)
    add_textbox(slide, 9.2, ry + 0.05, 1.85, 0.32, meaning,
                font_size=Pt(8), color=TEXT_DARK)
    add_textbox(slide, 11.15, ry + 0.05, 1.75, 0.32, use,
                font_size=Pt(8), color=TEXT_DARK, wrap=True)
    ry += 0.42

# Readiness gate callout
callout_box(slide, 7.1, 4.42, 5.9, 0.75,
    "Readiness Gate: The solver will not run unless all contracts are approved, SAP-validated, "
    "have non-stale open positions (refreshed within 30 minutes), and all external API data is "
    "within its polling interval.",
    border_color=PURPLE_DARK, bg=PURPLE_LIGHT, text_color=PURPLE_DARK)

# Constraint bridge examples
add_textbox(slide, 0.4, 5.18, 12, 0.28, "WHAT THE CONSTRAINT BRIDGE MEANS IN PRACTICE",
            font_size=Pt(9), bold=True, color=BLUE)
cb_items = [
    ("Quantity Floor", "A long-term supply contract requires minimum 5,000 MT/month to St. Louis. The Bridge sets a lower bound — the solver cannot recommend less."),
    ("Penalty Margin", "A contract carries $15/ton shortfall penalty if below 80% of nominated volume. The Bridge reduces effective sell price by probability-weighted penalty."),
    ("Price Override", "A fixed-price contract locks delivery at $400/ton regardless of spot. The Bridge replaces sell_stl with $400 for the covered tonnage."),
    ("Delivery Window", "Contract specifies delivery between the 15th and 25th. If the window is approaching, the Bridge elevates the minimum allocation to ensure on-time commitment."),
]
cx = 0.4
for title, desc in cb_items:
    small_box(slide, cx, 5.5, 3.05, 1.4, title, desc)
    cx += 3.23

footer(slide, "Section 4: Contract Ingestion & SAP Integration")

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — SECTION 5: HOW THE SOLVER WORKS
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "5", "How the Solver Finds the Optimal Solution",
    "The solver uses the HiGHS linear programming engine to find the exact allocation of product "
    "to routes that optimises the chosen objective while satisfying every constraint simultaneously.")

# LP explanation
add_textbox(slide, 0.4, 2.0, 5.8, 0.3, "LINEAR PROGRAMMING EXAMPLE",
            font_size=Pt(9), bold=True, color=BLUE)

lp_text = (
    "Maximise: profit(Don→StL)×x₁ + profit(Don→Mem)×x₂ + profit(Geis→StL)×x₃ + profit(Geis→Mem)×x₄\n\n"
    "Subject to:\n"
    "  x₁ + x₂ ≤ inv_don          (Donaldsonville inventory)\n"
    "  x₃ + x₄ ≤ inv_geis         (Geismar inventory)\n"
    "  (x₁+x₂+x₃+x₄)/eff_barge ≤ barge_count   (fleet capacity)\n"
    "  (x₁+x₂+x₃+x₄) × nola_buy ≤ working_cap  (capital)\n"
    "  x₁ + x₃ ≥ contract_floor_stl             (contract minimum)\n"
    "  x₂ = 0  if mem_outage = true              (dock outage)\n"
    "  xᵢ ≥ 0  for all routes"
)
add_rect(slide, 0.4, 2.3, 6.0, 2.3, fill=LIGHT_BLUE, line_color=BLUE, line_width=Pt(2))
tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.38), Inches(5.75), Inches(2.15))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = lp_text
run.font.size = Pt(8.5)
run.font.color.rgb = DARK_BLUE

# Monte Carlo section
add_textbox(slide, 6.8, 2.0, 6, 0.3, "MONTE CARLO: QUANTIFYING RISK",
            font_size=Pt(9), bold=True, color=BLUE)
add_textbox(slide, 6.8, 2.3, 6.1, 0.55,
    "The platform perturbs all 20 variables across 1,000 scenarios (e.g., river stage ±2 ft, "
    "NOLA buy ±$10/ton) and solves each independently to produce a full probability distribution.",
    font_size=Pt(9.5), color=TEXT_MID, wrap=True)

# Signal badges
signals = [
    ("STRONG GO", RGBColor(0x1b, 0x5e, 0x20), RGBColor(0xc8, 0xe6, 0xc9), RGBColor(0xa5, 0xd6, 0xa7),
     "5th percentile profit positive. Even worst-case scenarios are profitable."),
    ("GO", RGBColor(0x2e, 0x7d, 0x32), RGBColor(0xdc, 0xed, 0xc8), RGBColor(0xc5, 0xe1, 0xa5),
     "25th percentile positive. Good probability of solid outcome."),
    ("CAUTIOUS", RGBColor(0xe6, 0x51, 0x00), RGBColor(0xff, 0xe0, 0xb2), RGBColor(0xff, 0xcc, 0x80),
     "Median positive but downside scenarios could be negative."),
    ("WEAK", RGBColor(0xbf, 0x36, 0x0c), RGBColor(0xff, 0xcc, 0xbc), RGBColor(0xff, 0xab, 0x91),
     "Median positive but narrow. High sensitivity to variable moves."),
    ("NO GO", RGBColor(0xb7, 0x1c, 0x1c), RGBColor(0xff, 0xcd, 0xd2), RGBColor(0xef, 0x9a, 0x9a),
     "Median negative or infeasible in most scenarios."),
]
sx = 6.8
for sig_name, fg, bg, border, desc in signals:
    add_rect(slide, sx, 2.92, 1.12, 0.65, fill=bg, line_color=border, line_width=Pt(1.5))
    add_textbox(slide, sx + 0.04, 2.96, 1.04, 0.28, sig_name,
                font_size=Pt(8.5), bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_textbox(slide, sx + 0.04, 3.22, 1.04, 0.35, desc,
                font_size=Pt(7), color=fg, wrap=True, align=PP_ALIGN.CENTER)
    sx += 1.23

# Shadow prices
add_textbox(slide, 0.4, 4.75, 12, 0.3, "SHADOW PRICES — KNOWING WHAT TO NEGOTIATE",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 0.4, 5.05, 12.4, 0.32, fill=NAVY)
for col_x, col_w, col_text in [(0.45, 3.5, "Constraint"), (3.95, 4.5, "Shadow Price Meaning"), (8.5, 4.2, "Trader Action")]:
    add_textbox(slide, col_x, 5.08, col_w, 0.26, col_text, font_size=Pt(9), bold=True, color=WHITE)

sp_rows = [
    ("Barge fleet is binding", "$X additional profit per extra barge available", "Charter a spot barge if shadow price exceeds charter rate"),
    ("Working capital is binding", "$Y additional profit per $1,000 of extra capital", "Assess whether short-term credit is worth finance cost"),
    ("Inventory at Don is binding", "$Z additional profit per extra ton at Donaldsonville", "Evaluate opportunistic spot purchase to build inventory"),
    ("Contract minimum is binding", "Negative shadow price (constraint is costly)", "Escalate to commercial — obligation is destroying margin"),
]
ry = 5.37
for i, (con, meaning, action) in enumerate(sp_rows):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    add_rect(slide, 0.4, ry, 12.4, 0.38, fill=bg)
    add_textbox(slide, 0.45, ry + 0.04, 3.4, 0.3, con,
                font_size=Pt(8.5), bold=True, color=DARK_BLUE, wrap=True)
    add_textbox(slide, 3.95, ry + 0.04, 4.4, 0.3, meaning,
                font_size=Pt(8.5), color=TEXT_DARK, wrap=True)
    add_textbox(slide, 8.5, ry + 0.04, 4.2, 0.3, action,
                font_size=Pt(8.5), color=TEXT_DARK, wrap=True)
    ry += 0.38

# Sensitivity
callout_box(slide, 0.4, 6.95, 12.4, 0.38,
    "Sensitivity Analysis: After Monte Carlo, the platform reports which variables most strongly "
    "correlate with profit outcomes — e.g., 'Freight rates explain 42% of variance in this trade's profitability.'")

footer(slide, "Section 5: How the Solver Finds the Optimal Solution")

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — SECTION 6: THE AUTO TRADER
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "6", "The Auto Trader",
    "The auto-runner continuously monitors all 20 live variables and automatically re-optimises "
    "whenever conditions move beyond a configured sensitivity threshold. It records every optimal "
    "outcome to an immutable audit trail.")

# Flow diagram
add_textbox(slide, 0.4, 2.0, 12, 0.3, "CONTINUOUS MONITORING LOOP",
            font_size=Pt(9), bold=True, color=BLUE)
flow_y = 2.3
flow_box(slide, 0.4, flow_y, 2.6, 0.7, "Live Variables", "Continuous API updates")
arrow(slide, 3.05, flow_y + 0.2)
flow_box(slide, 3.4, flow_y, 3.0, 0.7, "Delta Detector", "|current − centre| ≥ threshold?",
         border=AMBER_DARK, bg=AMBER_LIGHT, fg=AMBER_DARK)
arrow(slide, 6.45, flow_y + 0.2)
flow_box(slide, 6.8, flow_y, 2.5, 0.7, "LP + Monte Carlo", "1,000 scenarios",
         border=NAVY, bg=NAVY, fg=STEEL_BLUE)
arrow(slide, 9.35, flow_y + 0.2)
flow_box(slide, 9.7, flow_y, 3.0, 0.7, "Optimal Outcome", "New centre point saved",
         border=GREEN_DARK, bg=GREEN_LIGHT, fg=GREEN_DARK)
add_textbox(slide, 7.5, 3.05, 5, 0.3,
            "↑ Also: trader notification + audit trail commit",
            font_size=Pt(8.5), color=TEXT_LIGHT, italic=True)

# Trigger thresholds table
add_textbox(slide, 0.4, 3.5, 6, 0.28, "EXAMPLE DELTA THRESHOLDS (NH3 Domestic Barge)",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 0.4, 3.78, 6.2, 0.3, fill=NAVY)
for col_x, col_w, col_text in [(0.45, 1.4, "Variable"), (1.9, 1.0, "Threshold"), (3.0, 3.5, "Why This Level")]:
    add_textbox(slide, col_x, 3.81, col_w, 0.24, col_text, font_size=Pt(8.5), bold=True, color=WHITE)
thresh_rows = [
    ("river_stage", "±0.5 ft", "Half a foot can shift barge capacity tier"),
    ("nola_buy", "±$2/ton", "NH3 buy price moves in $2 increments; smaller moves are noise"),
    ("lock_hrs", "±2 hrs", "Two hours of delay affects transit economics and timing"),
    ("fr_don_stl", "±$3/ton", "$3 represents a decision-relevant freight shift"),
    ("barge_count", "±1 barge", "Any change in available fleet directly changes max allocation"),
    ("inv_don", "±500 tons", "~One-third of a barge load — material supply swing"),
]
ty = 4.08
for i, (var, thr, why) in enumerate(thresh_rows):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    add_rect(slide, 0.4, ty, 6.2, 0.36, fill=bg)
    add_textbox(slide, 0.45, ty + 0.05, 1.35, 0.26, var, font_size=Pt(8.5), bold=True, color=DARK_BLUE)
    add_textbox(slide, 1.9, ty + 0.05, 0.95, 0.26, thr, font_size=Pt(8.5), color=TEXT_DARK)
    add_textbox(slide, 3.0, ty + 0.05, 3.45, 0.26, why, font_size=Pt(8), color=TEXT_MID, wrap=True)
    ty += 0.36

# Post-solve steps on right side
add_textbox(slide, 7.0, 3.5, 5.8, 0.28, "WHAT HAPPENS AFTER A RE-SOLVE",
            font_size=Pt(9), bold=True, color=BLUE)
post_steps = [
    (NAVY, "1", "Optimal Outcome Recorded",
     "Route allocations, profit, margins, MC distribution saved to ETS, Postgres, and SQLite."),
    (BLUE, "2", "Immutable Chain Commit",
     "Solve payload is hashed (SHA-256), signed (ECDSA), encrypted (AES-256-GCM), and written to BSV blockchain as OP_RETURN."),
    (PURPLE_DARK, "3", "AI Analyst Explanation",
     "Platform generates plain-English explanation of why the auto-runner triggered and what changed in the recommendation."),
    (AMBER_DARK, "4", "Trader Notification",
     "If profit change exceeds the trader's configured threshold, an alert is sent via their preferred channels."),
]
py = 3.8
for color, num, title, desc in post_steps:
    add_rect(slide, 7.0, py, 0.37, 0.37, fill=color)
    add_textbox(slide, 7.0, py + 0.05, 0.37, 0.27, num,
                font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.45, py + 0.02, 5.3, 0.22, title,
                font_size=Pt(9.5), bold=True, color=TEXT_DARK)
    add_textbox(slide, 7.45, py + 0.24, 5.3, 0.25, desc,
                font_size=Pt(8.5), color=TEXT_MID, wrap=True)
    py += 0.58

# Cooldown callout
callout_box(slide, 0.4, 6.25, 12.4, 0.6,
    "Minimum Cooldown: The auto-runner enforces a minimum 5-minute interval between solves (configurable). "
    "A fallback scheduled solve runs every 60 minutes regardless of delta, ensuring the recommendation "
    "stays fresh even in quiet markets.",
    border_color=GREEN_DARK, bg=GREEN_LIGHT, text_color=RGBColor(0x1b, 0x5e, 0x20))

footer(slide, "Section 6: The Auto Trader")

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — SECTION 7: TRADER NOTIFICATIONS
# ══════════════════════════════════════════════════════════════════
slide = add_slide()
section_header(slide, "7", "Trader Notifications & Threshold Alerts",
    "The notification system ensures that when conditions shift enough to change the optimal decision "
    "materially, the right trader is told — on their preferred channel — with enough context to act immediately.")

# Three trigger conditions
add_textbox(slide, 0.4, 2.0, 12, 0.28, "WHAT TRIGGERS A NOTIFICATION — ALL THREE CONDITIONS MUST BE TRUE",
            font_size=Pt(9), bold=True, color=BLUE)
cond_items = [
    ("Condition 1 — Auto-solve triggered",
     "The auto-runner has solved because a variable moved beyond its delta threshold. "
     "A scheduled fallback solve alone does not trigger a notification."),
    ("Condition 2 — Profit delta exceeds threshold",
     "The profit change from last solve exceeds the trader's notify_threshold_profit "
     "(default $5,000). Small fluctuations don't generate noise."),
    ("Condition 3 — Cooldown elapsed",
     "At least notify_cooldown_minutes (default 30 min) have passed since the last "
     "notification. Back-to-back alerts are suppressed."),
]
cx = 0.4
for title, desc in cond_items:
    small_box(slide, cx, 2.3, 3.97, 1.1, title, desc)
    cx += 4.17

# Notification example
add_textbox(slide, 0.4, 3.5, 12, 0.28, "SAMPLE NOTIFICATION",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 0.4, 3.78, 7.7, 2.45, fill=LIGHT_BLUE, line_color=BLUE, line_width=Pt(2))
notif_text = (
    "Subject: NH3 Domestic — Auto Solve: $142,500 (▲ $18,200)\n\n"
    "Hi Sarah,\n\n"
    "The Trammo NH3 Domestic Trading Desk has triggered an auto-solve.\n\n"
    "Current Expected Profit: $142,500\n"
    "Change from Last Solve: ▲ $18,200\n"
    "Triggered by: river_stage (was 11.8 ft → now 18.3 ft, threshold ±0.5 ft)\n\n"
    "Monte Carlo (1,000 scenarios):\n"
    "  Mean: $138,900   VaR₅: $62,400   P95: $210,500\n"
    "  Signal: STRONG GO — 97% of scenarios profitable\n\n"
    "Analyst Note: River stage recovered to 18.3 ft, restoring full barge capacity "
    "(1,500 MT/barge). The binding constraint shifted from fleet capacity back to "
    "working capital. Recommend committing to Don→StL at full allocation before "
    "freight rates respond to the gauge reading."
)
tb = slide.shapes.add_textbox(Inches(0.55), Inches(3.86), Inches(7.45), Inches(2.3))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = notif_text
run.font.size = Pt(8.5)
run.font.color.rgb = DARK_BLUE

# Channels and settings on right
add_textbox(slide, 8.3, 3.5, 5, 0.28, "NOTIFICATION CHANNELS",
            font_size=Pt(9), bold=True, color=BLUE)
channels = [
    ("Email", "Full detail; useful when away from desk or for record-keeping"),
    ("Slack", "Real-time team visibility; multiple traders can see the same alert"),
    ("MS Teams", "Preferred for organisations running Teams as primary collaboration platform"),
]
cy = 3.78
for ch, desc in channels:
    small_box(slide, 8.3, cy, 4.9, 0.68, ch, desc)
    cy += 0.73

add_textbox(slide, 8.3, 5.0, 5, 0.28, "TRADER-CONFIGURABLE SETTINGS",
            font_size=Pt(9), bold=True, color=BLUE)
add_rect(slide, 8.3, 5.28, 4.9, 0.28, fill=NAVY)
add_textbox(slide, 8.35, 5.31, 2.5, 0.22, "Setting", font_size=Pt(8.5), bold=True, color=WHITE)
add_textbox(slide, 10.9, 5.31, 0.8, 0.22, "Default", font_size=Pt(8.5), bold=True, color=WHITE)
settings = [
    ("notify_threshold_profit", "$5,000", "Minimum profit change to trigger notification"),
    ("notify_cooldown_minutes", "30 min", "Minimum time between notifications"),
    ("notifications_paused", "false", "One-click pause; auto-runner continues solving"),
]
sy = 5.56
for i, (setting, default, effect) in enumerate(settings):
    bg = LIGHT_GREY if i % 2 == 1 else WHITE
    add_rect(slide, 8.3, sy, 4.9, 0.38, fill=bg)
    add_textbox(slide, 8.35, sy + 0.05, 2.45, 0.28, setting,
                font_size=Pt(8), bold=True, color=DARK_BLUE)
    add_textbox(slide, 10.85, sy + 0.05, 0.9, 0.28, default,
                font_size=Pt(8), color=TEXT_DARK)
    sy += 0.38

# Trader workflow
add_textbox(slide, 0.4, 6.3, 12, 0.28, "THE TRADER'S OPTIMISATION LOOP",
            font_size=Pt(9), bold=True, color=BLUE)
loop_steps = [
    ("1", "Receive notification", "Alert arrives with optimal outcome, what changed, and analyst's plain-English interpretation."),
    ("2", "Open Trader Mode", "Use dashboard sliders to test variations. Solver responds in under a second."),
    ("3", "Review shadow prices", "Shadow prices show where negotiations deliver most value."),
    ("4", "Commit the trade", "Execute in the market. Auto-runner picks up new centre point and resumes monitoring."),
]
lx = 0.4
for num, title, desc in loop_steps:
    colors = [NAVY, BLUE, GREEN_DARK, AMBER_DARK]
    color = colors[int(num)-1]
    add_rect(slide, lx, 6.58, 0.35, 0.35, fill=color)
    add_textbox(slide, lx, 6.63, 0.35, 0.25, num,
                font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, lx + 0.4, 6.6, 2.6, 0.22, title,
                font_size=Pt(9), bold=True, color=TEXT_DARK)
    add_textbox(slide, lx + 0.4, 6.82, 2.6, 0.45, desc,
                font_size=Pt(8), color=TEXT_MID, wrap=True)
    lx += 3.23

footer(slide, "Section 7: Trader Notifications & Threshold Alerts")

# ── Save ─────────────────────────────────────────────────────────
output_path = "/home/user/trading_desk/executive_overview.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
